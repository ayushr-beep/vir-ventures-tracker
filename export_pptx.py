"""
Vir Ventures — 9-slide Boardroom PPTX with Vendor Analysis slide.
"""
import io,numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches,Pt,Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NV=RGBColor(0x0F,0x16,0x29); NV2=RGBColor(0x1E,0x2A,0x45)
OR=RGBColor(0xE8,0x61,0x1A); OR2=RGBColor(0xF4,0x87,0x4B)
WH=RGBColor(0xFF,0xFF,0xFF); LT=RGBColor(0xF8,0xF9,0xFB)
LG=RGBColor(0xE5,0xE7,0xEB); GY=RGBColor(0x6B,0x72,0x80)
GR=RGBColor(0x16,0xA3,0x4A); RD=RGBColor(0xDC,0x26,0x26)
AM=RGBColor(0xD9,0x77,0x06); BL=RGBColor(0x25,0x63,0xEB)
DK=RGBColor(0x11,0x18,0x27)
VCLRS={"CRDI":OR,"PNCL":BL,"PROH":AM}
VNAMES={"CRDI":"Gaming & Entertainment","PNCL":"The Pencil Grip","PROH":"Prohitter Sports"}
W=Inches(13.33); H=Inches(7.5)

YES_VALS={"yes","yes, need discount"}
def is_yes(v): return str(v).strip().lower() in YES_VALS

def _r(slide,x,y,w,h,fill,line=None):
    s=slide.shapes.add_shape(1,x,y,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line: s.line.color.rgb=line
    else: s.line.fill.background()
    return s

def _t(slide,text,x,y,w,h,sz=11,bold=False,color=WH,align=PP_ALIGN.LEFT,italic=False):
    tb=slide.shapes.add_textbox(x,y,w,h)
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r2=p.add_run(); r2.text=str(text)
    r2.font.size=Pt(sz); r2.font.bold=bold
    r2.font.italic=italic; r2.font.color.rgb=color
    r2.font.name="Calibri"
    return tb

def _bg(slide,color):
    fill=slide.background.fill; fill.solid(); fill.fore_color.rgb=color

def _hdr(slide,title,sub,week):
    _r(slide,Emu(0),Emu(0),W,Inches(.75),NV)
    _t(slide,title.upper(),Inches(.45),Inches(.12),Inches(9),Inches(.5),sz=10,bold=True,color=OR)
    _t(slide,week,W-Inches(3.2),Inches(.12),Inches(2.9),Inches(.5),sz=10,color=LG,align=PP_ALIGN.RIGHT)
    _t(slide,sub,Inches(.45),Inches(.82),Inches(11),Inches(.3),sz=11,color=GY)

def _ftr(slide):
    _r(slide,Emu(0),H-Inches(.35),W,Inches(.35),NV)
    _t(slide,"VIR VENTURES  ·  BOARDROOM ANALYST REPORT  ·  CONFIDENTIAL",
       Inches(.4),H-Inches(.3),W-Inches(.8),Inches(.25),sz=7,bold=True,color=RGBColor(0x3A,0x42,0x60))

def _kpi(slide,x,y,w,h,label,val,sub,acc=OR,vc=None):
    _r(slide,x,y,w,h,WH)
    _r(slide,x,y,w,Inches(.05),acc)
    _t(slide,label,x+Inches(.15),y+Inches(.1),w-Inches(.2),Inches(.2),sz=7,bold=True,color=GY)
    _t(slide,val, x+Inches(.15),y+Inches(.32),w-Inches(.2),Inches(.62),sz=26,bold=True,color=vc or DK)
    _t(slide,sub, x+Inches(.15),y+Inches(.95),w-Inches(.2),Inches(.2),sz=9,color=GY)

# ─── SLIDE 1: COVER ──────────────────────────────────────────────────────────
def s_cover(prs,d):
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,NV)
    _r(s,W-Inches(3.6),Emu(0),Inches(3.6),Inches(2.4),OR)
    _r(s,W-Inches(2.0),Emu(0),Inches(2.0),H,NV2)
    _t(s,"VIR VENTURES",Inches(.7),Inches(1.4),Inches(7),Inches(.5),sz=13,bold=True,color=OR)
    _t(s,"Boardroom Analyst\nPerformance Report",Inches(.7),Inches(2.1),Inches(8.5),Inches(2.2),sz=44,bold=True,color=WH)
    _t(s,d['week_label'],Inches(.7),Inches(4.5),Inches(7),Inches(.45),sz=16,color=OR2)
    _t(s,f"Report Date: {d['report_date']}  ·  SKUs: {d['total']}  ·  Approval Rate: {d['rate']}%",Inches(.7),Inches(5.1),Inches(9),Inches(.3),sz=11,color=GY)
    _r(s,Emu(0),H-Inches(.5),W-Inches(2.0),Inches(.5),OR)
    _t(s,"ANALYST PERFORMANCE TRACKER  ·  WEEKLY BOARDROOM REVIEW",Inches(.4),H-Inches(.44),W-Inches(2.4),Inches(.38),sz=9,bold=True,color=WH)
    s.notes_slide.notes_text_frame.text="Cover slide. Open with week overview and introduce the three key sections: Vendor Analysis, ASIN Intelligence, and Analyst Performance."

# ─── SLIDE 2: EXECUTIVE SUMMARY ──────────────────────────────────────────────
def s_exec(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"Executive Summary","At-a-glance performance snapshot for management",d['week_label']); _ftr(s)

    cw=Inches(2.8); ch=Inches(1.42); gap=Inches(.2); ky=Inches(1.2)
    kpis=[("CATALOGUE SIZE",str(d['total']),"Total SKUs reviewed",OR,DK),
          ("APPROVED",str(d['approved']),f"{d['rate']}% approval rate",GR,GR),
          ("REJECTED",str(d['rejected']),"Negative/Low margin",RD,RD),
          ("APPROVAL RATE",f"{d['rate']}%",f"of {d['total']} SKUs",BL,BL)]
    for i,(lbl,val,sub,acc,vc) in enumerate(kpis):
        _kpi(s,Inches(.4)+i*(cw+gap),ky,cw,ch,lbl,val,sub,acc,vc)

    yl=df[df['is_yes']]
    top_a=df.groupby('Analyst')['is_yes'].mean().idxmax() if not yl.empty else "—"
    top_ar=round(df.groupby('Analyst')['is_yes'].mean().max()*100,1) if not yl.empty else 0
    top_brand=yl.groupby('Brand')['ASIN'].count().idxmax() if not yl.empty else "—"
    top_brand_n=int(yl.groupby('Brand')['ASIN'].count().max()) if not yl.empty else 0
    top_asin=yl.sort_values('Last month sold',ascending=False).iloc[0] if not yl.empty else None
    avg_m=round(yl['margin_pct'].mean(),1) if not yl.empty else 0
    vp_count=df['Vendor Prefix'].nunique() if 'Vendor Prefix' in df.columns else "—"

    _r(s,Inches(.4),Inches(2.82),Inches(12.5),Inches(3.85),WH)
    _t(s,"KEY MANAGEMENT INSIGHTS",Inches(.6),Inches(2.92),Inches(8),Inches(.28),sz=8,bold=True,color=OR)
    bullets=[
        f"This week {d['total']} SKUs reviewed across {df['Brand'].nunique()} brands and {vp_count} vendor groups. Overall approval rate: {d['rate']}%.",
        f"{top_a} is top analyst — {top_ar}% approval rate. Strong demand-score selection with low FBA competition picks.",
        f"Average Buy Box margin on approved SKUs: {avg_m}% — well above breakeven across all vendors.",
        f"{top_brand} leads brand performance with {top_brand_n} approved SKUs — immediate PO action recommended.",
        f"Top ASIN: {top_asin['ASIN'] if top_asin is not None else '—'} ({top_asin['Brand'] if top_asin is not None else '—'}) — {int(top_asin['Last month sold']) if top_asin is not None else 0} units/month.",
        f"{d['rejected']} SKUs rejected (all PROH vendor) — negative/low margin. Vendor renegotiation required.",
    ]
    by=Inches(3.35)
    for b in bullets:
        _r(s,Inches(.55),by,Inches(.055),Inches(.2),OR)
        _t(s,b,Inches(.72),by-Inches(.02),Inches(11.1),Inches(.36),sz=11,color=DK)
        by+=Inches(.43)
    s.notes_slide.notes_text_frame.text="Walk through KPIs first. Then deliver each insight as a talking point. Assign actions before leaving slide."

# ─── SLIDE 3: VENDOR ANALYSIS ────────────────────────────────────────────────
def s_vendor(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"Vendor Analysis","Performance breakdown by vendor group — CRDI · PNCL · PROH",d['week_label']); _ftr(s)

    if 'Vendor Prefix' not in df.columns:
        _t(s,"No Vendor Prefix column in data.",Inches(2),Inches(3),Inches(9),Inches(1),sz=14,color=GY)
        return

    vp=df.groupby('Vendor Prefix').agg(
        Total=('ASIN','count'),Approved=('is_yes','sum'),
        Rejected=('is_yes',lambda x:(~x).sum()),
        Brands=('Brand','nunique'),Units=('Last month sold','sum'),
        AvgM=('margin_pct','mean'),AvgRank=('Rank','mean'),
        AvgDem=('Total Demand Score','mean'),AvgComp=('Total Competition Score','mean'),
        AvgMarSc=('Total Margin Score','mean'),
    ).reset_index().round(1)
    vp['Rate']=(vp['Approved']/vp['Total']*100).round(1)

    n=len(vp); cw2=Inches(12.3/n); sy=Inches(1.2); gap2=Inches(.18)

    for i,(_,row) in enumerate(vp.iterrows()):
        pfx=row['Vendor Prefix']
        clr=VCLRS.get(pfx,OR)
        rate=row['Rate']
        rate_clr=GR if rate>=80 else (AM if rate>=50 else RD)
        cx=Inches(.4)+i*(cw2+gap2)

        _r(s,cx,sy,cw2,Inches(5.6),WH)
        _r(s,cx,sy,cw2,Inches(.07),clr)

        _t(s,pfx,cx+Inches(.18),sy+Inches(.12),cw2-Inches(.25),Inches(.3),sz=11,bold=True,color=clr)
        _t(s,VNAMES.get(pfx,pfx),cx+Inches(.18),sy+Inches(.45),cw2-Inches(.25),Inches(.35),sz=10,color=DK)

        # Rate big number
        _t(s,f"{rate}%",cx+Inches(.18),sy+Inches(.88),cw2-Inches(.25),Inches(.72),sz=32,bold=True,color=rate_clr)
        _t(s,"Approval Rate",cx+Inches(.18),sy+Inches(1.6),cw2-Inches(.25),Inches(.22),sz=8,color=GY)

        stats=[("Total SKUs",str(int(row['Total']))),
               ("Approved",str(int(row['Approved']))),
               ("Rejected",str(int(row['Rejected']))),
               ("Units/Month",str(int(row['Units']))),
               ("Avg Margin %",f"{row['AvgM']:.1f}%"),
               ("Avg BSR Rank",f"#{row['AvgRank']:,.0f}"),
               ("Brands",str(int(row['Brands']))),
               ("Demand Score",str(row['AvgDem']))]
        sy2=sy+Inches(1.9)
        for lbl2,val2 in stats:
            _r(s,cx+Inches(.15),sy2,cw2-Inches(.25),Inches(.4),LT)
            _t(s,lbl2,cx+Inches(.2),sy2+Inches(.02),cw2-Inches(.35),Inches(.18),sz=7,color=GY)
            _t(s,val2,cx+Inches(.2),sy2+Inches(.2),cw2-Inches(.35),Inches(.18),sz=10,bold=True,
               color=GR if lbl2=="Approved" else (RD if lbl2=="Rejected" else (OR if "Margin" in lbl2 else DK)))
            sy2+=Inches(.46)

    # Key vendor insight
    best_v=vp.loc[vp['Rate'].idxmax(),'Vendor Prefix']
    worst_v=vp.loc[vp['Rate'].idxmin(),'Vendor Prefix']
    best_units=int(vp.loc[vp['Units'].idxmax(),'Units'])
    best_units_v=vp.loc[vp['Units'].idxmax(),'Vendor Prefix']
    _r(s,Inches(.4),Inches(7.0),Inches(12.5),Inches(.2),OR)
    note=f"Best vendor: {best_v} ({vp.loc[vp['Vendor Prefix']==best_v,'Rate'].values[0]:.0f}% approval)  ·  Most demand: {best_units_v} ({best_units:,} units/mo)  ·  Needs attention: {worst_v}"
    _t(s,note,Inches(.5),Inches(7.0),Inches(12.0),Inches(.2),sz=9,bold=True,color=WH)
    s.notes_slide.notes_text_frame.text="Compare the three vendor cards side by side. CRDI = gaming, high-value SKUs. PNCL = stationery, 100% approval. PROH = problematic, all rejections come from here."

# ─── SLIDE 4: TOP ASINs ──────────────────────────────────────────────────────
def s_asins(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"ASIN Intelligence — Top Approved SKUs","Ranked by monthly demand velocity",d['week_label']); _ftr(s)

    yl=df[df['is_yes']].sort_values('Last month sold',ascending=False).head(10)
    _r(s,Inches(.4),Inches(1.2),Inches(12.5),Inches(.4),NV)
    hdrs=[("ASIN",1.5),("Vendor",0.85),("Brand",2.0),("Analyst",.9),("Rec.",.85),
          ("Units/Mo",.85),("Net $",.7),("BB $",.7),("Margin%",.8),("Rank",.95),("Score",.65),("BB Holder",1.3)]
    cx=Inches(.4)
    for hdr,cw3 in hdrs:
        _t(s,hdr,cx+Inches(.05),Inches(1.25),Inches(cw3-.08),Inches(.3),sz=7,bold=True,color=WH,align=PP_ALIGN.CENTER)
        cx+=Inches(cw3)

    for ri,(_,row) in enumerate(yl.iterrows()):
        ry=Inches(1.62)+ri*Inches(.51)
        bg=WH if ri%2==0 else LT
        _r(s,Inches(.4),ry,Inches(12.5),Inches(.49),bg)
        if ri==0: _r(s,Inches(.4),ry,Inches(.06),Inches(.49),OR)
        cells=[
            (str(row.get('ASIN',''))[:12],1.5,DK,ri==0),
            (str(row.get('Vendor Prefix','')),.85,VCLRS.get(str(row.get('Vendor Prefix','')),GY),True),
            (str(row.get('Brand',''))[:18],2.0,DK,ri==0),
            (str(row.get('Analyst','')), .9,GY,False),
            (str(row.get('Recommended',''))[:12],.85,OR if ri==0 else GY,False),
            (str(int(row.get('Last month sold',0))),.85,GR,True),
            (f"${row.get('Net price',0):.2f}",.7,DK,False),
            (f"${row.get('BB Price',0):.2f}",.7,DK,False),
            (f"{row.get('margin_pct',0):.1f}%",.8,OR,True),
            (f"#{int(row.get('Rank',0)):,}",.95,DK,False),
            (str(row.get('Total Product Score','')),.65,DK,False),
            (str(row.get('BuyBoxSellerName',''))[:14],1.3,GY,False),
        ]
        tx=Inches(.4)
        for val,cw3,clr,bld in cells:
            _t(s,val,tx+Inches(.05),ry+Inches(.12),Inches(cw3-.08),Inches(.26),sz=8.5,bold=bld,color=clr,align=PP_ALIGN.CENTER)
            tx+=Inches(cw3)
    s.notes_slide.notes_text_frame.text="Focus on rows highlighted in orange (rank 1). Green = high units/month. Top 5 rows should be actioned for immediate PO this week."

# ─── SLIDE 5: ANALYST SCORECARDS ─────────────────────────────────────────────
def s_analysts(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"Analyst Performance Scorecards","Individual quality, throughput, and scoring",d['week_label']); _ftr(s)

    astats=df.groupby('Analyst').agg(
        Total=('ASIN','count'),Approved=('is_yes','sum'),
        Rejected=('is_yes',lambda x:(~x).sum()),
        Rate=('is_yes',lambda x:round(x.mean()*100,1)),
        AvgM=('margin_pct','mean'),Units=('Last month sold','sum'),
        Rank=('Rank','mean'),Net=('Net price','mean'),BB=('BB Price','mean'),
        Dem=('Total Demand Score','mean'),Comp=('Total Competition Score','mean'),
        MarSc=('Total Margin Score','mean')).reset_index().round(1)

    n=len(astats); cw2=Inches(12.3/n); sy=Inches(1.2); gap2=Inches(.18)
    for i,(_,row) in enumerate(astats.iterrows()):
        rate=row['Rate']; rc2=GR if rate>=80 else (AM if rate>=50 else RD)
        cx=Inches(.4)+i*(cw2+gap2)
        _r(s,cx,sy,cw2,Inches(5.7),WH)
        _r(s,cx,sy,cw2,Inches(.07),rc2)
        _t(s,row['Analyst'],cx+Inches(.18),sy+Inches(.12),cw2-Inches(.25),Inches(.38),sz=16,bold=True,color=DK)
        _t(s,f"{rate}% APPROVAL",cx+Inches(.18),sy+Inches(.54),cw2-Inches(.25),Inches(.28),sz=10,bold=True,color=rc2)
        stats2=[("Total Audited",str(int(row['Total']))),
                ("Approved",str(int(row['Approved']))),
                ("Rejected",str(int(row['Rejected']))),
                ("Avg Margin %",f"{row['AvgM']:.1f}%"),
                ("Avg BSR Rank",f"#{row['Rank']:,.0f}"),
                ("Units/Month",str(int(row['Units']))),
                ("Avg Net $",f"${row['Net']:.2f}"),
                ("Avg BB $",f"${row['BB']:.2f}"),
                ("Demand Score",str(row['Dem'])),
                ("Comp. Score",str(row['Comp'])),
                ("Margin Score",str(row['MarSc']))]
        sy2=sy+Inches(.95)
        for lbl2,val2 in stats2:
            _r(s,cx+Inches(.15),sy2,cw2-Inches(.25),Inches(.4),LT)
            _t(s,lbl2,cx+Inches(.2),sy2+Inches(.02),cw2-Inches(.35),Inches(.17),sz=7,color=GY)
            _t(s,val2,cx+Inches(.2),sy2+Inches(.2),cw2-Inches(.35),Inches(.18),sz=10,bold=True,
               color=GR if lbl2=="Approved" else (RD if lbl2=="Rejected" else (OR if "Margin" in lbl2 else DK)))
            sy2+=Inches(.44)
    s.notes_slide.notes_text_frame.text="Green = 80%+ rate. Amber = 50-79%. Red = below 50%. Highlight top analyst and discuss workload balance. 100% rate analysts had clean vendor batches."

# ─── SLIDE 6: BRAND ANALYSIS ─────────────────────────────────────────────────
def s_brands(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"Brand Intelligence","Approved SKU count, margin, and demand by brand",d['week_label']); _ftr(s)

    bs=df[df['is_yes']].groupby('Brand').agg(
        Approved=('ASIN','count'),Units=('Last month sold','sum'),
        AvgM=('margin_pct','mean'),Rank=('Rank','mean'),
    ).reset_index().sort_values('Approved',ascending=False).head(10).round(1)
    mx=bs['Approved'].max()

    tw=Inches(7.5); lw=Inches(2.0); bh=Inches(.42); gap3=Inches(.1)
    sx=Inches(.4); sy=Inches(1.2)
    _t(s,"BRAND",sx,sy-Inches(.32),lw,Inches(.25),sz=8,bold=True,color=GY)
    _t(s,"APPROVED SKUs  →",sx+lw+Inches(.1),sy-Inches(.32),tw,Inches(.25),sz=8,bold=True,color=GY)

    for i,(_,row) in enumerate(bs.iterrows()):
        by=sy+i*(bh+gap3)
        pct=row['Approved']/mx if mx else 0
        fw2=max(Inches(.15),tw*pct)
        clr=OR if i==0 else (OR2 if i<3 else RGBColor(0xF9,0xC0,0x8E))
        _t(s,row['Brand'],sx,by+Inches(.09),lw,bh,sz=9,bold=(i==0),color=DK,align=PP_ALIGN.RIGHT)
        _r(s,sx+lw+Inches(.15),by,tw,bh,LG)
        _r(s,sx+lw+Inches(.15),by,fw2,bh,clr)
        _t(s,str(int(row['Approved'])),sx+lw+Inches(.15)+fw2+Inches(.08),by+Inches(.09),Inches(.5),bh,sz=10,bold=(i==0),color=DK)

    # Right metrics panel
    tx2=Inches(10.1); ty2=Inches(1.2)
    _t(s,"BRAND METRICS",tx2,ty2-Inches(.35),Inches(2.8),Inches(.28),sz=8,bold=True,color=OR)
    _r(s,tx2,ty2,Inches(2.8),Inches(.38),NV)
    for hdr2,cx_off,w2 in [("Brand",0,1.05),("Units",1.1,.85),("Margin",1.98,.82)]:
        _t(s,hdr2,tx2+Inches(cx_off)+Inches(.05),ty2+Inches(.08),Inches(w2),Inches(.22),
           sz=8,bold=True,color=WH,align=PP_ALIGN.CENTER)
    for ri2,(_,row) in enumerate(bs.iterrows()):
        ry2=ty2+Inches(.38)+ri2*Inches(.47)
        bg2=WH if ri2%2==0 else LT
        _r(s,tx2,ry2,Inches(2.8),Inches(.45),bg2)
        for val2,cx_off2,w2,clr2,bld2 in [
            (str(row['Brand'])[:12],0,1.05,DK,ri2==0),
            (str(int(row['Units'])),1.1,.85,GR,False),
            (f"{row['AvgM']:.0f}%",1.98,.82,OR,True)]:
            _t(s,val2,tx2+Inches(cx_off2)+Inches(.05),ry2+Inches(.12),Inches(w2),Inches(.22),
               sz=9,bold=bld2,color=clr2,align=PP_ALIGN.CENTER)
    s.notes_slide.notes_text_frame.text="Top 3 brands should drive purchase orders this week. Bar length = relative SKU share. Margin % and units together determine sourcing priority."

# ─── SLIDE 7: REJECTION ANALYSIS ─────────────────────────────────────────────
def s_rejected(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"Rejection Analysis — Root Cause","All rejected ASINs with margin breakdown and next actions",d['week_label']); _ftr(s)

    rej=df[~df['is_yes']]
    neg=int(rej['Remarks '].str.contains('Negative',na=False).sum()) if 'Remarks ' in rej.columns else 0
    low=int(rej['Remarks '].str.contains('Low',na=False).sum()) if 'Remarks ' in rej.columns else 0
    oth=len(rej)-neg-low

    for i,(lbl,val,clr) in enumerate([("NEGATIVE MARGIN",str(neg),RD),("LOW MARGIN",str(low),AM),("OTHER",str(oth),GY)]):
        _kpi(s,Inches(.4)+i*Inches(3.8),Inches(1.2),Inches(3.5),Inches(1.4),lbl,val,"PROH vendor",clr,clr)

    _r(s,Inches(.4),Inches(2.8),Inches(12.5),Inches(.4),NV)
    hdrs2=[("ASIN",1.5),("Vendor",.85),("Brand",1.6),("Analyst",.9),("Net $",.85),
           ("BB $",.85),("Breakeven",.95),("Diff SP",.85),("Score",.7),("Reason",2.4)]
    cx=Inches(.4)
    for hdr2,cw3 in hdrs2:
        _t(s,hdr2,cx+Inches(.05),Inches(2.85),Inches(cw3-.08),Inches(.28),sz=7,bold=True,color=WH,align=PP_ALIGN.CENTER)
        cx+=Inches(cw3)

    for ri,(_,row) in enumerate(rej.iterrows()):
        ry=Inches(3.22)+ri*Inches(.48)
        bg=RGBColor(0xFF,0xF1,0xF2) if ri%2==0 else WH
        _r(s,Inches(.4),ry,Inches(12.5),Inches(.46),bg)
        cells=[
            (str(row.get('ASIN',''))[:14],1.5,DK,False),
            (str(row.get('Vendor Prefix','')), .85,VCLRS.get(str(row.get('Vendor Prefix','')),GY),True),
            (str(row.get('Brand',''))[:16],1.6,DK,True),
            (str(row.get('Analyst','')), .9,GY,False),
            (f"${row.get('Net price',0):.2f}",.85,DK,False),
            (f"${row.get('BB Price',0):.2f}",.85,DK,False),
            (f"${row.get('Breakeven',0):.2f}",.95,RD,True),
            (f"${row.get('Difference from SP',0):.2f}",.85,RD,True),
            (str(row.get('Total Product Score','')),.7,GY,False),
            (str(row.get('Remarks ',''))[:28],2.4,RD,False),
        ]
        tx=Inches(.4)
        for val,cw3,clr,bld in cells:
            _t(s,val,tx+Inches(.05),ry+Inches(.11),Inches(cw3-.08),Inches(.26),sz=8,bold=bld,color=clr,align=PP_ALIGN.CENTER)
            tx+=Inches(cw3)
    s.notes_slide.notes_text_frame.text="All 8 rejections are PROH vendor. 4 negative margin, 4 low margin. Action: contact vendor for price concession before next PO."

# ─── SLIDE 8: SOURCING RECOMMENDATIONS ───────────────────────────────────────
def s_sourcing(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,LT)
    _hdr(s,"Sourcing Recommendations — Action Plan","Priority-ranked actions for sourcing and vendor management teams",d['week_label']); _ftr(s)

    yl=df[df['is_yes']]
    top5=yl.sort_values('Last month sold',ascending=False).head(5)
    top_brand=yl.groupby('Brand')['ASIN'].count().idxmax() if not yl.empty else "—"
    top_brand_m=round(yl.groupby('Brand')['margin_pct'].mean().get(top_brand,0),1) if not yl.empty else 0
    top_brand_n=int(yl.groupby('Brand')['ASIN'].count().max()) if not yl.empty else 0

    prios=[
        (OR,"P1 — IMMEDIATE PO ACTION",
         "Action these ASINs this week (highest demand velocity):\n"+
         "  ·  ".join([f"{r['ASIN']} ({r['Brand']}, {int(r['Last month sold'])} units/mo)" for _,r in top5.iterrows()])),
        (BL,f"P2 — VENDOR FOCUS: {top_brand}",
         f"{top_brand_n} approved SKUs · Avg margin {top_brand_m}% · Expand PO allocation and negotiate volume pricing for next cycle."),
        (AM,"P3 — VENDOR RENEGOTIATION: PROH",
         f"8 ASINs rejected (negative/low margin). Request 15-20% cost reduction. If unresolved, reduce PROH PO allocation next quarter."),
        (GR,"P4 — GROWTH OPPORTUNITY",
         "Nintendo (7 SKUs) and Ubisoft (4 SKUs) show strong FBA competition gaps (<5 sellers). Increase catalogue coverage and prioritise restocking."),
    ]
    py=Inches(1.2)
    for clr,title,body in prios:
        _r(s,Inches(.4),py,Inches(12.5),Inches(1.2),WH)
        _r(s,Inches(.4),py,Inches(.1),Inches(1.2),clr)
        _t(s,title,Inches(.62),py+Inches(.1),Inches(11.6),Inches(.28),sz=9,bold=True,color=clr)
        _t(s,body,Inches(.62),py+Inches(.42),Inches(11.6),Inches(.7),sz=10,color=DK)
        py+=Inches(1.35)
    s.notes_slide.notes_text_frame.text="Assign owners before leaving this slide. P1=Sourcing. P2=Vendor Mgr. P3=Finance+Procurement. P4=Catalogue."

# ─── SLIDE 9: NEXT STEPS ─────────────────────────────────────────────────────
def s_close(prs,d):
    df=d['df']
    s=prs.slides.add_slide(prs.slide_layouts[6]); _bg(s,NV)
    _r(s,Emu(0),Emu(0),Inches(4.6),H,OR)
    _t(s,"NEXT\nSTEPS",Inches(.55),Inches(1.8),Inches(3.7),Inches(2.8),sz=52,bold=True,color=WH)
    _t(s,d['week_label'],Inches(.55),Inches(4.8),Inches(4.0),Inches(.4),sz=10,color=WH)
    _t(s,"virventures.com",Inches(.55),H-Inches(.9),Inches(4.0),Inches(.35),sz=10,color=RGBColor(0xF9,0xC0,0x8E))

    yl=df[df['is_yes']]
    top_brand=yl.groupby('Brand')['ASIN'].count().idxmax() if not yl.empty else "—"
    steps=[("01","Sourcing Team","Issue POs for top 5 demand ASINs immediately."),
           ("02","Vendor Management",f"Expand {top_brand} allocation — negotiate volume pricing."),
           ("03","Finance & Procurement","PROH vendor: target 15-20% price reduction. No orders until resolved."),
           ("04","Analyst Team","Maintain weekly cadence. Target 75%+ approval rate all analysts."),
           ("05","Catalogue Team","Expand Nintendo + Ubisoft SKU coverage — low FBA competition gap."),
           ("06","Management","Review next week's mastersheet against this week's benchmarks.")]
    sy2=Inches(1.2)
    for num,owner,action in steps:
        _r(s,Inches(5.1),sy2,Inches(7.8),Inches(.82),NV2)
        _t(s,num,Inches(5.2),sy2+Inches(.12),Inches(.5),Inches(.58),sz=18,bold=True,color=OR)
        _t(s,owner.upper(),Inches(5.82),sy2+Inches(.08),Inches(6.8),Inches(.28),sz=8,bold=True,color=OR)
        _t(s,action,Inches(5.82),sy2+Inches(.32),Inches(6.8),Inches(.44),sz=10,color=WH)
        sy2+=Inches(.95)
    s.notes_slide.notes_text_frame.text="Assign each action item to a named owner. Set follow-up date for PROH vendor outcome. Close with next week's timeline."

# ─── MAIN ────────────────────────────────────────────────────────────────────
def generate_pptx(data:dict)->bytes:
    prs=Presentation()
    prs.slide_width=W; prs.slide_height=H
    s_cover(prs,data)
    s_exec(prs,data)
    s_vendor(prs,data)
    s_asins(prs,data)
    s_analysts(prs,data)
    s_brands(prs,data)
    s_rejected(prs,data)
    s_sourcing(prs,data)
    s_close(prs,data)
    buf=io.BytesIO(); prs.save(buf)
    return buf.getvalue()
